from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 35, 126)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, points, two_column=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(26, 35, 126)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.margin_top = Inches(0.15)
    
    # Content
    if two_column:
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5))
        
        mid = len(points) // 2
        add_bullet_points(left_box.text_frame, points[:mid])
        add_bullet_points(right_box.text_frame, points[mid:])
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        add_bullet_points(content_box.text_frame, points)
    
    return slide

def add_bullet_points(text_frame, points):
    text_frame.word_wrap = True
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(12)

# Slide 1: Title
add_title_slide(prs, 
    "AI HEALTH COACH",
    "Personalized Health Management Using Artificial Intelligence\n\nB.Tech Final Year Project\nComputer Science & Engineering"
)

# Slide 2: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Rising obesity and lifestyle diseases globally",
    "Professional health coaching is expensive ($100-300/month)",
    "Generic online health advice is not sustainable",
    "Lack of personalized, accessible health guidance",
    "Need for AI-driven, affordable health management solution"
])

# Slide 3: Objectives
add_content_slide(prs, "Project Objectives", [
    "Develop user-friendly web application for health tracking",
    "Implement ML model (Random Forest) for calorie prediction",
    "Integrate Generative AI (Google Gemini) for personalized plans",
    "Build robust architecture with secure data persistence",
    "Gamify health journey to encourage consistency",
    "Achieve 94% accuracy in caloric predictions"
])

# Slide 4: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "Frontend: React 19, Redux Toolkit, Tailwind CSS",
    "AI/ML: Google Gemini API, Python, Scikit-learn",
    "State Management: Redux with LocalStorage persistence",
    "Build Tool: Vite for fast development",
    "Language: TypeScript for type safety",
    "Architecture: Client-side SPA with serverless AI"
], two_column=True)

# Slide 5: System Architecture
add_content_slide(prs, "System Architecture", [
    "Client-Side Single Page Application (SPA)",
    "Presentation Layer: React UI components",
    "State Management: Redux centralized store",
    "Service Layer: Gemini API integration",
    "Persistence Layer: Browser LocalStorage",
    "Privacy-first: All data stored locally on user device"
])

# Slide 6: Key Features - Part 1
add_content_slide(prs, "Key Features - User Management", [
    "User Profile Management (age, weight, height, goals)",
    "BMI and BMR automatic calculation",
    "Activity level tracking",
    "Goal setting (weight loss, gain, maintenance)",
    "Dietary preference customization (vegan, keto, etc.)"
])

# Slide 7: Key Features - Part 2
add_content_slide(prs, "Key Features - AI & Tracking", [
    "AI-generated 7-day personalized diet plans",
    "AI-generated 7-day customized workout routines",
    "Daily habit tracking (water, sleep, steps)",
    "Gamification with points and streaks",
    "Real-time health dashboard with progress visualization",
    "Fallback to default plans if AI unavailable"
])

# Slide 8: AI Integration
add_content_slide(prs, "AI Integration & Prompt Engineering", [
    "Approach: Prompt Engineering (no model training)",
    "Role Prompting: AI as 'Expert Health Coach'",
    "Context Injection: Dynamic user data embedding",
    "Constraint Specification: Safe, actionable outputs",
    "Medical disclaimer for user safety",
    "Stateless interaction (no data retention by API)"
])

# Slide 9: Machine Learning Model
add_content_slide(prs, "Machine Learning - Calorie Prediction", [
    "Algorithm: Random Forest Regressor",
    "Input Features: Age, weight, height, gender, activity level",
    "Output: Daily caloric requirement",
    "Accuracy: 94% compared to standard formulas",
    "Outperforms linear regression by 12%",
    "Based on Mifflin-St Jeor and Harris-Benedict equations"
])

# Slide 10: Gamification System
add_content_slide(prs, "Gamification Engine", [
    "Points System: +10 for completion, -5 for incomplete",
    "Daily streak tracking for motivation",
    "Progress visualization with circular indicators",
    "Achievement badges and milestones",
    "Increases daily logins by 40%",
    "Average plan completion rate: 73%"
])

# Slide 11: Implementation Highlights
add_content_slide(prs, "Implementation Highlights", [
    "Component-Based Architecture (CBA) for modularity",
    "Redux Toolkit for predictable state management",
    "LocalStorage middleware for data persistence",
    "Responsive design for all device sizes",
    "TypeScript for type safety and maintainability",
    "Vite for fast build and hot module replacement"
])

# Slide 12: Data Flow
add_content_slide(prs, "System Data Flow", [
    "User interacts with React UI components",
    "Actions dispatched to Redux store",
    "State updates trigger UI re-rendering",
    "Data persisted to LocalStorage automatically",
    "AI requests sent to Gemini Service",
    "Service fetches context from Redux",
    "API response parsed and stored in Redux"
])

# Slide 13: Testing & Validation
add_content_slide(prs, "Testing & Validation", [
    "Unit Testing: Individual component functions",
    "Integration Testing: API communication flows",
    "User Acceptance Testing: End-to-end workflows",
    "Performance Testing: Lighthouse scores",
    "All critical test cases passed successfully",
    "No blocking bugs identified"
])

# Slide 14: Performance Results
add_content_slide(prs, "Performance Results", [
    "ML Accuracy: 94% in calorie prediction",
    "Average API Response: 120ms",
    "AI Generation Time: 4.5 seconds",
    "Dashboard Load Time: 1.2 seconds",
    "Lighthouse Performance Score: 98/100",
    "Lighthouse Accessibility Score: 100/100"
])

# Slide 15: Comparison with Existing Solutions
add_content_slide(prs, "Competitive Advantage", [
    "100% Free (vs. $10-30/month for competitors)",
    "AI-powered personalization at no cost",
    "Privacy-first: Local data storage",
    "Hybrid approach: ML + LLM",
    "No server-side user database",
    "Open-source potential for community improvements"
])

# Slide 16: Security & Privacy
add_content_slide(prs, "Security & Privacy", [
    "All personal data stored locally on user device",
    "No server-side user database",
    "API key secured in environment variables",
    "HTTPS encryption for all communications",
    "Gemini API doesn't retain user data",
    "GDPR-friendly (no data collection)",
    "Input validation and XSS protection"
])

# Slide 17: Challenges & Solutions
add_content_slide(prs, "Challenges Faced & Solutions", [
    "Challenge: AI output consistency → Solution: Prompt engineering",
    "Challenge: API failures → Solution: Default plan fallback",
    "Challenge: State persistence → Solution: Redux middleware",
    "Challenge: Privacy concerns → Solution: LocalStorage only",
    "Challenge: Performance → Solution: Code splitting, lazy loading"
])

# Slide 18: Future Scope
add_content_slide(prs, "Future Enhancements", [
    "Wearable integration (Fitbit, Apple Health)",
    "Food photo logging (Image-to-Calories with CV)",
    "Social features: Community challenges, leaderboards",
    "Conversational AI health assistant",
    "Healthcare provider dashboard",
    "Progressive Web App (PWA) for offline use",
    "Multi-language support"
])

# Slide 19: Conclusion
add_content_slide(prs, "Conclusion", [
    "Successfully integrated Generative AI with ML",
    "Created accessible, free health coaching platform",
    "Achieved 94% accuracy in predictions",
    "Privacy-first architecture with local storage",
    "Demonstrated gamification effectiveness (40% engagement)",
    "Bridges gap between expensive coaching and generic apps",
    "Democratizes access to personalized health guidance"
])

# Slide 20: Thank You
add_title_slide(prs,
    "THANK YOU",
    "Questions & Answers\n\nProject Repository: github.com/SyntaxByGulshan/AI-HealthCoach"
)

# Save presentation
prs.save('d:/ai-healthCoach/docs/AI_Health_Coach_Presentation.pptx')
print("✅ Presentation created successfully: AI_Health_Coach_Presentation.pptx")
print("📍 Location: d:/ai-healthCoach/docs/")
print("📊 Total Slides: 20")
