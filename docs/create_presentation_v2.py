from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 35, 126)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, points, two_column=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(26, 35, 126)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.margin_top = Inches(0.15)
    
    # Content
    if two_column:
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5.5))
        
        mid = len(points) // 2
        add_bullet_points(left_box.text_frame, points[:mid])
        add_bullet_points(right_box.text_frame, points[mid:])
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        add_bullet_points(content_box.text_frame, points)
    
    return slide

def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(26, 35, 126)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.margin_top = Inches(0.15)
    
    # Add image if it exists
    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.8)
        height = Inches(5)
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        # Add placeholder text if image doesn't exist
        text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        text_frame = text_box.text_frame
        text_frame.text = f"[Diagram: {title}]"
        text_frame.paragraphs[0].font.size = Pt(28)
        text_frame.paragraphs[0].font.italic = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_bullet_points(text_frame, points):
    text_frame.word_wrap = True
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(12)

# Slide 1: Title
add_title_slide(prs, 
    "AI HEALTH COACH",
    "Personalized Health Management Using Artificial Intelligence\n\nB.Tech Final Year Project\nComputer Science & Engineering"
)

# Slide 2: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Rising obesity and lifestyle diseases globally",
    "Professional health coaching is expensive ($100-300/month)",
    "Generic online health advice is not sustainable",
    "Lack of personalized, accessible health guidance",
    "Need for AI-driven, affordable health management solution"
])

# Slide 3: Objectives
add_content_slide(prs, "Project Objectives", [
    "Develop user-friendly web application for health tracking",
    "Implement ML model (Random Forest) for calorie prediction",
    "Integrate Generative AI (Google Gemini) for personalized plans",
    "Build robust architecture with secure data persistence",
    "Gamify health journey to encourage consistency",
    "Achieve 94% accuracy in caloric predictions"
])

# Slide 4: Development Approach
add_content_slide(prs, "Development Approach", [
    "Component-Based Architecture (CBA)",
    "Agile Development Model for iterative refinement",
    "Modular code with independent, reusable components",
    "Unidirectional data flow with Redux",
    "Declarative UI with React",
    "Iterative prompt engineering for AI accuracy"
])

# Slide 5: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "Frontend: React 19 (Vite), Redux Toolkit, Tailwind CSS 4",
    "AI/ML: Google Gemini API, Python 3.11, Scikit-learn",
    "State Management: Redux with LocalStorage persistence",
    "Language: TypeScript for type safety",
    "Architecture: Client-side SPA with serverless AI",
    "Build Tool: Vite for fast development experience"
], two_column=True)

# Slide 6: System Architecture Layers
add_content_slide(prs, "System Architecture - Layered Approach", [
    "Presentation Layer: React UI components and hooks",
    "State Management Layer: Redux centralized store",
    "Service Layer: Gemini API integration service",
    "Persistence Layer: LocalStorage for browser durability",
    "Cloud Layer: Google Gemini API (serverless)",
    "Privacy-first: All user data stored locally"
])

# Slide 7: System Modules
add_content_slide(prs, "Key System Modules", [
    "User Profile Management: Demographics, BMI, BMR calculation",
    "Daily Habits & Gamification: Points, streaks, progress tracking",
    "AI Plan Generation: Diet and workout personalization",
    "Footsteps Tracker: Physical activity integration",
    "State Management: Redux slices for each module",
    "Persistence: LocalStorage sync middleware"
])

# Slide 8: Use Case Diagram
add_image_slide(prs, "Use Case Diagram", "d:/ai-healthCoach/docs/diagrams/use_case_diagram.png")

# Slide 9: Class Diagram
add_image_slide(prs, "Class Diagram", "d:/ai-healthCoach/docs/diagrams/class_diagram.png")

# Slide 10: Activity Diagram
add_image_slide(prs, "Activity Diagram - AI Plan Generation", "d:/ai-healthCoach/docs/diagrams/activity_diagram.png")

# Slide 11: Sequence Diagram
add_image_slide(prs, "Sequence Diagram - AI Plan Generation", "d:/ai-healthCoach/docs/diagrams/sequence_diagram.png")

# Slide 12: Component Diagram
add_image_slide(prs, "Component Diagram", "d:/ai-healthCoach/docs/diagrams/component_diagram.png")

# Slide 13: Deployment Diagram
add_image_slide(prs, "Deployment Diagram", "d:/ai-healthCoach/docs/diagrams/deployment_diagram.png")

# Slide 14: AI Integration & Prompt Engineering
add_content_slide(prs, "AI Integration & Prompt Engineering", [
    "Approach: Prompt Engineering (no model training)",
    "Role Prompting: AI as 'Expert Health Coach'",
    "Context Injection: Dynamic user data embedding",
    "Constraint Specification: Safe, actionable outputs",
    "Fallback mechanism: Default plans if API fails",
    "Verification: Output validation for safety"
])

# Slide 15: Machine Learning Model
add_content_slide(prs, "Machine Learning - Calorie Prediction", [
    "Algorithm: Random Forest Regressor",
    "Input: Age, weight, height, gender, activity level",
    "Output: Daily caloric requirement",
    "Accuracy: 94% vs. standard formulas",
    "Outperforms linear regression by 12%",
    "Based on Mifflin-St Jeor & Harris-Benedict equations"
])

# Slide 16: Testing & Results
add_content_slide(prs, "Testing & Performance Results", [
    "ML Accuracy: 94% in calorie prediction",
    "Average API Response: 120ms",
    "AI Generation Time: 4.5 seconds",
    "Dashboard Load: 1.2 seconds",
    "Lighthouse Performance: 98/100",
    "Lighthouse Accessibility: 100/100",
    "All critical test cases passed"
])

# Slide 17: Security & Privacy
add_content_slide(prs, "Security & Privacy Considerations", [
    "Local Data Storage: All data on user's device",
    "API Security: Environment variables for keys",
    "Stateless AI: Gemini doesn't retain user data",
    "HTTPS encryption for all communications",
    "Input validation and XSS protection",
    "GDPR-friendly (no data collection)",
    "No server-side user database"
])

# Slide 18: Deployment Plan
add_content_slide(prs, "Deployment Strategy", [
    "Build: Vite compiles TypeScript and bundles assets",
    "Hosting: Netlify, Vercel, or GitHub Pages",
    "CI/CD: Automatic builds and deployments",
    "Environment: Production and development configs",
    "Monitoring: Error tracking and analytics",
    "Scalability: Serverless architecture for easy scaling"
])

# Slide 19: Future Scope
add_content_slide(prs, "Future Enhancements", [
    "Wearable integration (Fitbit, Apple Health)",
    "Food photo logging (Image-to-Calories with CV)",
    "Social features: Community challenges, leaderboards",
    "Conversational AI health assistant",
    "Healthcare provider dashboard",
    "Progressive Web App (PWA) for offline use",
    "Multi-language support"
])

# Slide 20: Conclusion
add_content_slide(prs, "Conclusion", [
    "Successfully integrated Generative AI with ML",
    "Created accessible, free health coaching platform",
    "Achieved 94% accuracy in predictions",
    "Privacy-first architecture with local storage",
    "Component-based, modular design",
    "Agile methodology enabled iterative refinement",
    "Bridges gap between expensive coaching and generic apps"
])

# Slide 21: Thank You
add_title_slide(prs,
    "THANK YOU",
    "Questions & Answers\n\nProject Repository: github.com/SyntaxByGulshan/AI-HealthCoach"
)

# Save presentation
prs.save('d:/ai-healthCoach/docs/AI_Health_Coach_Presentation.pptx')
print("✅ Presentation updated successfully: AI_Health_Coach_Presentation.pptx")
print("📍 Location: d:/ai-healthCoach/docs/")
print("📊 Total Slides: 21")
print("🎨 Includes: UML diagrams, architecture, development approach")
